"""Build Memory Psychology midterm PPT (20 slides).

Topic: Foreign Language Grammar Learning — Declarative → Procedural Memory.
Sources: Anderson (2000) Ch.1/6/7/8/9/10, Fitts & Posner (1967), DeKeyser (1997, 2007),
Ullman (2001, 2004), Newell & Rosenbloom (1981).
"""
from __future__ import annotations

import os
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# ---------------- Style Tokens ----------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

C_BG = RGBColor(0xFF, 0xFF, 0xFF)
C_TITLE = RGBColor(0x00, 0x00, 0x00)        # black
C_TEXT = RGBColor(0x1A, 0x1A, 0x1A)         # near-black
C_MUTED = RGBColor(0x66, 0x66, 0x66)        # mid gray
C_ACCENT = RGBColor(0x33, 0x33, 0x33)       # dark gray
C_LIGHT = RGBColor(0xF2, 0xF2, 0xF2)        # light gray bg
C_LINE = RGBColor(0x99, 0x99, 0x99)         # gray line

# Stage/system distinctions removed — use labels for hierarchy
C_COG = RGBColor(0x00, 0x00, 0x00)
C_ASSOC = RGBColor(0x00, 0x00, 0x00)
C_AUTO = RGBColor(0x00, 0x00, 0x00)

C_DECL = RGBColor(0x00, 0x00, 0x00)
C_PROC = RGBColor(0x00, 0x00, 0x00)

KOR_FONT = "Malgun Gothic"
ENG_FONT = "Calibri"

OUT_DIR = Path.home() / "Desktop"
OUT_DIR.mkdir(parents=True, exist_ok=True)
PPT_PATH = OUT_DIR / "Memory_Psychology_L2_Grammar.pptx"
ASSET_DIR = Path("/tmp/mempsych_assets")
ASSET_DIR.mkdir(exist_ok=True)


# ---------------- Helpers ----------------
def set_white_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = C_BG


def add_textbox(slide, x, y, w, h, text, *, size=22, bold=False, color=C_TEXT,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if font:
        run.font.name = font
    return tb


def add_multiline(slide, x, y, w, h, lines, *, size=22, bold=False, color=C_TEXT,
                  align=PP_ALIGN.LEFT, line_space=1.2, anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, opts_dict) tuples or plain strings."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, item in enumerate(lines):
        if isinstance(item, tuple):
            text, opts = item
        else:
            text, opts = item, {}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get("align", align)
        p.line_spacing = line_space
        p.space_after = Pt(opts.get("space_after", 6))
        run = p.add_run()
        run.text = text
        run.font.size = Pt(opts.get("size", size))
        run.font.bold = opts.get("bold", bold)
        run.font.color.rgb = opts.get("color", color)
    return tb


def add_title_bar(slide, title_main, subtitle=None):
    """Standard slide title at top."""
    add_textbox(slide, Inches(0.7), Inches(0.4), Inches(12.0), Inches(0.8),
                title_main, size=32, bold=True, color=C_TITLE)
    if subtitle:
        add_textbox(slide, Inches(0.7), Inches(1.05), Inches(12.0), Inches(0.45),
                    subtitle, size=18, color=C_MUTED)
    # underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.7),
                                   Inches(1.55) if subtitle else Inches(1.25),
                                   Inches(1.2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = C_ACCENT
    line.line.fill.background()


def add_box(slide, x, y, w, h, fill=None, line=None, line_w=1.0, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    box = slide.shapes.add_shape(shape_type, x, y, w, h)
    if fill is None:
        box.fill.background()
    else:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    if line is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = line
        box.line.width = Pt(line_w)
    box.shadow.inherit = False
    return box


def add_arrow(slide, x1, y1, x2, y2, color=C_ACCENT, weight=3.0):
    conn = slide.shapes.add_connector(2, x1, y1, x2, y2)  # straight connector with arrow
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    # add arrow end
    line_elem = conn.line._get_or_add_ln()
    from pptx.oxml.ns import qn
    from lxml import etree
    tail = etree.SubElement(line_elem, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("h", "med")
    return conn


def add_page_number(slide, n, total):
    add_textbox(slide, Inches(12.4), Inches(7.05), Inches(0.8), Inches(0.3),
                f"{n} / {total}", size=11, color=C_MUTED, align=PP_ALIGN.RIGHT)


def add_footer(slide, text="Memory Psychology · Foreign Language Grammar Learning"):
    add_textbox(slide, Inches(0.7), Inches(7.05), Inches(11.0), Inches(0.3),
                text, size=11, color=C_MUTED)


# ---------------- Power-law curve image ----------------
def make_power_curve(path):
    fig, ax = plt.subplots(figsize=(7, 4.0), dpi=180)
    n = np.linspace(1, 1000, 500)
    rt = 2200 * (n ** -0.4) + 380
    ax.plot(n, rt, color="black", linewidth=2)
    for x, label in [(20, "Cognitive"),
                      (200, "Associative"),
                      (700, "Autonomous")]:
        y = 2200 * (x ** -0.4) + 380
        ax.scatter([x], [y], s=80, color="black", zorder=5)
        ax.annotate(label, (x, y), xytext=(15, 14), textcoords="offset points",
                    fontsize=11, color="black")
    ax.set_xlabel("Practice trials (N)", fontsize=11)
    ax.set_ylabel("Reaction time (ms)", fontsize=11)
    ax.set_title("Power Law of Practice", fontsize=13, color="black", pad=10)
    ax.grid(True, linestyle=":", alpha=0.4, color="gray")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)


# ---------------- Build presentation ----------------
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]
TOTAL = 20


def new_slide():
    s = prs.slides.add_slide(BLANK)
    set_white_bg(s)
    return s


# ============= SLIDE 1 — Title =============
s = new_slide()
# accent strip on left
strip = add_box(s, Inches(0), Inches(0), Inches(0.35), SLIDE_H, fill=C_ACCENT)

add_textbox(s, Inches(0.9), Inches(0.7), Inches(11.0), Inches(0.5),
            "Memory Psychology  ·  Midterm Alternative Assignment",
            size=16, color=C_MUTED)

add_textbox(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(1.0),
            "Foreign Language Grammar Learning",
            size=44, bold=True, color=C_TITLE)

add_textbox(s, Inches(0.9), Inches(2.55), Inches(11.5), Inches(0.6),
            "Declarative → Procedural Memory의 점진적 전이",
            size=24, color=C_TEXT)

# Center diagram: Declarative → Procedural
y_d = Inches(4.0)
b1 = add_box(s, Inches(2.0), y_d, Inches(3.8), Inches(1.4),
             fill=C_LIGHT, line=C_DECL, line_w=2.0, rounded=True)
add_textbox(s, Inches(2.0), Inches(4.15), Inches(3.8), Inches(0.5),
            "DECLARATIVE", size=18, bold=True, color=C_DECL, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(2.0), Inches(4.55), Inches(3.8), Inches(0.5),
            "Knowing THAT", size=22, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(2.0), Inches(4.95), Inches(3.8), Inches(0.4),
            "explicit · conscious", size=14, color=C_MUTED, align=PP_ALIGN.CENTER)

add_arrow(s, Inches(6.0), Inches(4.7), Inches(7.5), Inches(4.7), color=C_ACCENT, weight=4)

b2 = add_box(s, Inches(7.7), y_d, Inches(3.8), Inches(1.4),
             fill=RGBColor(0xF2, 0xF2, 0xF2), line=C_PROC, line_w=2.0, rounded=True)
add_textbox(s, Inches(7.7), Inches(4.15), Inches(3.8), Inches(0.5),
            "PROCEDURAL", size=18, bold=True, color=C_PROC, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(7.7), Inches(4.55), Inches(3.8), Inches(0.5),
            "Knowing HOW", size=22, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(7.7), Inches(4.95), Inches(3.8), Inches(0.4),
            "implicit · automatic", size=14, color=C_MUTED, align=PP_ALIGN.CENTER)

# bottom info
add_textbox(s, Inches(0.9), Inches(6.4), Inches(11.0), Inches(0.4),
            "영어영문학과 4학년  ·  May 2026", size=14, color=C_MUTED)
add_page_number(s, 1, TOTAL)


# ============= SLIDE 2 — Overview =============
s = new_slide()
add_title_bar(s, "Overview", "오늘의 발표 흐름  ·  4 Parts")

parts = [
    ("01", "Memory Systems", "Declarative vs Procedural — 두 종류의 기억"),
    ("02", "Three Stages", "Cognitive → Associative → Autonomous"),
    ("03", "Mechanisms", "Production rules · Knowledge compilation · Power law"),
    ("04", "Application", "L2 grammar — DeKeyser & Ullman"),
]
for i, (num, head, sub) in enumerate(parts):
    y = Inches(2.1 + i * 1.15)
    # number badge
    badge = add_box(s, Inches(0.9), y, Inches(1.0), Inches(0.95),
                    fill=C_ACCENT, rounded=True)
    add_textbox(s, Inches(0.9), y + Inches(0.18), Inches(1.0), Inches(0.6),
                num, size=28, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(2.2), y + Inches(0.05), Inches(10.0), Inches(0.55),
                head, size=26, bold=True, color=C_TITLE)
    add_textbox(s, Inches(2.2), y + Inches(0.55), Inches(10.0), Inches(0.45),
                sub, size=16, color=C_MUTED)

add_footer(s); add_page_number(s, 2, TOTAL)


# ============= SLIDE 3 — Core Question =============
s = new_slide()
add_title_bar(s, "Core Question", "본 발표가 답할 질문")

add_textbox(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.0),
            "“Is grammar just memorization?”",
            size=44, bold=True, color=C_TITLE, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.9), Inches(3.15), Inches(11.5), Inches(0.6),
            "외국어 문법 학습은 단순 규칙 암기인가?",
            size=22, color=C_TEXT, align=PP_ALIGN.CENTER)

# Answer box
ans = add_box(s, Inches(2.5), Inches(4.4), Inches(8.3), Inches(1.6),
              fill=C_LIGHT, line=C_ACCENT, line_w=2.0, rounded=True)
add_textbox(s, Inches(2.5), Inches(4.55), Inches(8.3), Inches(0.5),
            "ANSWER", size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(2.5), Inches(4.95), Inches(8.3), Inches(0.6),
            "No — it is a shift in memory systems.",
            size=26, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(2.5), Inches(5.50), Inches(8.3), Inches(0.5),
            "기억 체계와 정보 처리 방식이 점차 변화하는 과정",
            size=18, color=C_MUTED, align=PP_ALIGN.CENTER)

add_footer(s); add_page_number(s, 3, TOTAL)


# ============= SLIDE 4 — Two Memory Systems =============
s = new_slide()
add_title_bar(s, "Two Memory Systems", "기억은 하나가 아니다")

# Left box — Declarative
lx = Inches(0.9); rx = Inches(7.0); bw = Inches(5.4); bh = Inches(4.8); by = Inches(2.0)
add_box(s, lx, by, bw, bh, fill=C_LIGHT, line=C_DECL, line_w=2.0, rounded=True)
add_textbox(s, lx, by + Inches(0.25), bw, Inches(0.5),
            "DECLARATIVE", size=16, bold=True, color=C_DECL, align=PP_ALIGN.CENTER)
add_textbox(s, lx, by + Inches(0.75), bw, Inches(0.6),
            "Knowing THAT", size=30, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)
decl_lines = [
    ("· 사실·사건에 대한 의식적 기억", {}),
    ("· Explicit · Episodic", {"color": C_MUTED}),
    ("· 언어로 진술 가능", {}),
    ("· Hippocampus 의존", {"color": C_MUTED}),
    ("· 빠른 습득 / 빠른 망각", {}),
]
add_multiline(s, lx + Inches(0.5), by + Inches(1.7), bw - Inches(0.9), Inches(3.0),
              decl_lines, size=20, line_space=1.4)

# Right box — Procedural
add_box(s, rx, by, bw, bh, fill=RGBColor(0xF2, 0xF2, 0xF2), line=C_PROC, line_w=2.0, rounded=True)
add_textbox(s, rx, by + Inches(0.25), bw, Inches(0.5),
            "PROCEDURAL", size=16, bold=True, color=C_PROC, align=PP_ALIGN.CENTER)
add_textbox(s, rx, by + Inches(0.75), bw, Inches(0.6),
            "Knowing HOW", size=30, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)
proc_lines = [
    ("· 기능·절차의 무의식적 기억", {}),
    ("· Implicit · Automatic", {"color": C_MUTED}),
    ("· 진술 어려움", {}),
    ("· Basal ganglia · Cerebellum", {"color": C_MUTED}),
    ("· 느린 습득 / 강한 보존", {}),
]
add_multiline(s, rx + Inches(0.5), by + Inches(1.7), bw - Inches(0.9), Inches(3.0),
              proc_lines, size=20, line_space=1.4)

add_footer(s); add_page_number(s, 4, TOTAL)


# ============= SLIDE 5 — Declarative Memory =============
s = new_slide()
add_title_bar(s, "Declarative Memory", "선언적 기억  ·  Knowing THAT")

add_textbox(s, Inches(0.9), Inches(1.95), Inches(11.5), Inches(0.6),
            "사실·사건에 대해 의식적으로 인출되는 기억",
            size=22, color=C_TEXT, align=PP_ALIGN.LEFT)
add_textbox(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(0.4),
            "Memories that we are explicitly aware of  —  Anderson (2000), Ch.8",
            size=14, color=C_MUTED)

# 3 feature cards
features = [
    ("Explicit", "의식적으로 인출", C_DECL),
    ("Episodic", "맥락·시간이 함께 저장\n(언제, 어디서, 누구와)", C_DECL),
    ("Verbalizable", "말로 설명 가능\n규칙·예문·일화", C_DECL),
]
for i, (head, body, col) in enumerate(features):
    x = Inches(0.9 + i * 4.0)
    y = Inches(3.4)
    add_box(s, x, y, Inches(3.6), Inches(2.0), fill=C_LIGHT, line=col, line_w=1.5, rounded=True)
    add_textbox(s, x, y + Inches(0.25), Inches(3.6), Inches(0.6),
                head, size=24, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_textbox(s, x + Inches(0.3), y + Inches(0.95), Inches(3.0), Inches(1.0),
                body, size=16, color=C_TEXT, align=PP_ALIGN.CENTER)

# Example box
ex_y = Inches(5.7)
add_box(s, Inches(0.9), ex_y, Inches(11.5), Inches(1.1), fill=RGBColor(0xF7, 0xF7, 0xF7),
        line=C_LINE, line_w=1.0, rounded=True)
add_textbox(s, Inches(1.1), ex_y + Inches(0.1), Inches(2.0), Inches(0.4),
            "예시 EXAMPLE", size=12, bold=True, color=C_ACCENT)
add_textbox(s, Inches(1.1), ex_y + Inches(0.45), Inches(11.0), Inches(0.6),
            '“3인칭 단수 현재형엔 동사에 -s를 붙인다”  →  교사의 설명을 그대로 떠올림',
            size=18, color=C_TEXT)

add_footer(s); add_page_number(s, 5, TOTAL)


# ============= SLIDE 6 — Procedural Memory =============
s = new_slide()
add_title_bar(s, "Procedural Memory", "절차적 기억  ·  Knowing HOW")

add_textbox(s, Inches(0.9), Inches(1.95), Inches(11.5), Inches(0.6),
            "기능·습관·행위 패턴에 대한 무의식적 기억",
            size=22, color=C_TEXT)
add_textbox(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(0.4),
            "Memory underlying skill performance  —  Anderson (2000), Ch.9",
            size=14, color=C_MUTED)

# 3 feature cards
features = [
    ("Implicit", "의식 없이 작동", C_PROC),
    ("Automatic", "주의 자원 거의 안 씀\n이중과제 가능", C_PROC),
    ("Pattern-based", "조건 → 행동 형태로\n패키지화된 절차", C_PROC),
]
for i, (head, body, col) in enumerate(features):
    x = Inches(0.9 + i * 4.0)
    y = Inches(3.4)
    add_box(s, x, y, Inches(3.6), Inches(2.0), fill=RGBColor(0xF2, 0xF2, 0xF2),
            line=col, line_w=1.5, rounded=True)
    add_textbox(s, x, y + Inches(0.25), Inches(3.6), Inches(0.6),
                head, size=24, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_textbox(s, x + Inches(0.3), y + Inches(0.95), Inches(3.0), Inches(1.0),
                body, size=16, color=C_TEXT, align=PP_ALIGN.CENTER)

# Evidence box
ev_y = Inches(5.7)
add_box(s, Inches(0.9), ev_y, Inches(11.5), Inches(1.1), fill=RGBColor(0xF7, 0xF7, 0xF7),
        line=C_LINE, line_w=1.0, rounded=True)
add_textbox(s, Inches(1.1), ev_y + Inches(0.1), Inches(2.5), Inches(0.4),
            "결정적 증거 EVIDENCE", size=12, bold=True, color=C_ACCENT)
add_textbox(s, Inches(1.1), ev_y + Inches(0.45), Inches(11.0), Inches(0.6),
            "환자 H.M. — 해마 손상 후 새 사실은 못 외워도, 거울 따라 그리기 같은 skill은 정상 학습",
            size=18, color=C_TEXT)

add_footer(s); add_page_number(s, 6, TOTAL)


# ============= SLIDE 7 — Three Stages =============
s = new_slide()
add_title_bar(s, "Three Stages of Skill Acquisition",
              "Fitts & Posner (1967)  ·  Anderson (1982)")

stages = [
    ("Cognitive", "인지 / 선언", C_COG, "Knowing the rule"),
    ("Associative", "연합", C_ASSOC, "Connecting rule\nto context"),
    ("Autonomous", "자율 / 절차", C_AUTO, "Using without\nthinking"),
]
y = Inches(2.6)
bw = Inches(3.4); bh = Inches(2.4)
gap = Inches(0.6)
total_w = bw * 3 + gap * 2
start_x = (SLIDE_W - total_w) / 2

for i, (head, sub, col, line) in enumerate(stages):
    x = start_x + (bw + gap) * i
    add_box(s, x, y, bw, bh, fill=C_BG, line=col, line_w=2.5, rounded=True)
    # color band
    add_box(s, x, y, bw, Inches(0.5), fill=col, line=col, line_w=2.5, rounded=True)
    add_textbox(s, x, y + Inches(0.05), bw, Inches(0.4),
                f"Stage {i+1}", size=14, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    add_textbox(s, x, y + Inches(0.65), bw, Inches(0.55),
                head, size=24, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_textbox(s, x, y + Inches(1.15), bw, Inches(0.4),
                sub, size=14, color=C_MUTED, align=PP_ALIGN.CENTER)
    add_textbox(s, x, y + Inches(1.6), bw, Inches(0.7),
                line, size=16, color=C_TEXT, align=PP_ALIGN.CENTER)
    if i < 2:
        ax1 = x + bw + Inches(0.05)
        ax2 = x + bw + gap - Inches(0.05)
        add_arrow(s, ax1, y + bh / 2, ax2, y + bh / 2, color=C_ACCENT, weight=4)

# bottom — what changes
ch_y = Inches(5.6)
add_textbox(s, Inches(0.9), ch_y, Inches(11.5), Inches(0.4),
            "What changes across stages",
            size=14, bold=True, color=C_ACCENT)
add_textbox(s, Inches(0.9), ch_y + Inches(0.4), Inches(11.5), Inches(1.0),
            "Speed ↑     Errors ↓     Working-memory load ↓     Conscious effort ↓",
            size=20, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)

add_footer(s); add_page_number(s, 7, TOTAL)


# ============= SLIDE 8 — Stage 1 Cognitive =============
s = new_slide()
# stage badge
add_box(s, Inches(0.7), Inches(0.4), Inches(1.4), Inches(0.55), fill=C_COG, rounded=True)
add_textbox(s, Inches(0.7), Inches(0.45), Inches(1.4), Inches(0.5),
            "STAGE 1", size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER)
add_textbox(s, Inches(2.3), Inches(0.4), Inches(11.0), Inches(0.7),
            "Cognitive Stage", size=32, bold=True, color=C_TITLE)
add_textbox(s, Inches(2.3), Inches(1.05), Inches(11.0), Inches(0.4),
            "인지 단계  ·  Declarative knowledge 의존", size=18, color=C_MUTED)
line = add_box(s, Inches(0.7), Inches(1.55), Inches(1.2), Inches(0.04), fill=C_COG)

key_lines = [
    ("• 규칙과 예문을 설명 형태로 이해·암기", {}),
    ("• 작업기억 부담이 큼  →  발화가 느리고 의식적", {}),
    ("• 오류가 많고 일관성 부족", {}),
    ("• 본질적으로 problem solving 활동", {"color": C_MUTED}),
]
add_multiline(s, Inches(0.9), Inches(1.95), Inches(7.0), Inches(3.5),
              key_lines, size=20, line_space=1.5)

# Right panel — student's inner monologue example
panel_x = Inches(8.2); panel_y = Inches(1.95); panel_w = Inches(4.2); panel_h = Inches(4.5)
add_box(s, panel_x, panel_y, panel_w, panel_h, fill=C_LIGHT, line=C_COG,
        line_w=1.5, rounded=True)
add_textbox(s, panel_x, panel_y + Inches(0.2), panel_w, Inches(0.4),
            "학습자의 내적 사고", size=12, bold=True, color=C_COG, align=PP_ALIGN.CENTER)
add_textbox(s, panel_x + Inches(0.25), panel_y + Inches(0.7), panel_w - Inches(0.5),
            Inches(0.6),
            'Q.  "She ___ to school."',
            size=18, bold=True, color=C_TEXT)
inner = [
    ('"주어가 she니까…', {"size": 16}),
    ('  3인칭 단수,', {"size": 16}),
    ('  현재시제니까…', {"size": 16}),
    ('  동사에 -s를 붙여야 해."', {"size": 16}),
    ("", {"size": 8}),
    ("→  She goes to school.", {"size": 18, "bold": True, "color": C_TITLE}),
    ("(약 5–7초 소요)", {"size": 12, "color": C_MUTED}),
]
add_multiline(s, panel_x + Inches(0.3), panel_y + Inches(1.4),
              panel_w - Inches(0.6), Inches(2.8), inner, line_space=1.25)

add_textbox(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.4),
            "Anderson (2000), Ch.9, p.310 — “A skill develops from the cognitive stage…”",
            size=11, color=C_MUTED)

add_footer(s); add_page_number(s, 8, TOTAL)


# ============= SLIDE 9 — Stage 2 Associative =============
s = new_slide()
add_box(s, Inches(0.7), Inches(0.4), Inches(1.4), Inches(0.55), fill=C_ASSOC, rounded=True)
add_textbox(s, Inches(0.7), Inches(0.45), Inches(1.4), Inches(0.5),
            "STAGE 2", size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER)
add_textbox(s, Inches(2.3), Inches(0.4), Inches(11.0), Inches(0.7),
            "Associative Stage", size=32, bold=True, color=C_TITLE)
add_textbox(s, Inches(2.3), Inches(1.05), Inches(11.0), Inches(0.4),
            "연합 단계  ·  Knowledge compilation 시작", size=18, color=C_MUTED)
add_box(s, Inches(0.7), Inches(1.55), Inches(1.2), Inches(0.04), fill=C_ASSOC)

key_lines = [
    ("• 규칙과 사용 맥락 사이의 연결이 강화", {}),
    ("• 설명문 → IF-THEN production 형태로 재구조화", {}),
    ("• 익숙한 표현은 빨라지지만, 새 구조엔 멈칫", {}),
    ("• Declarative → Procedural 전이가 진행되는 구간", {"color": C_MUTED}),
]
add_multiline(s, Inches(0.9), Inches(1.95), Inches(7.0), Inches(3.5),
              key_lines, size=20, line_space=1.5)

# Right panel — example
panel_x = Inches(8.2); panel_y = Inches(1.95); panel_w = Inches(4.2); panel_h = Inches(4.5)
add_box(s, panel_x, panel_y, panel_w, panel_h, fill=RGBColor(0xF2, 0xF2, 0xF2),
        line=C_ASSOC, line_w=1.5, rounded=True)
add_textbox(s, panel_x, panel_y + Inches(0.2), panel_w, Inches(0.4),
            "L2 학습자 예시", size=12, bold=True, color=C_ASSOC, align=PP_ALIGN.CENTER)
ex_lines = [
    ('익숙한 표현 →', {"size": 14, "color": C_MUTED}),
    ('"He goes / She likes"', {"size": 18, "bold": True}),
    ('  즉각 발화 ✓', {"size": 14, "color": C_AUTO}),
    ("", {"size": 6}),
    ('새 동사 →', {"size": 14, "color": C_MUTED}),
    ('"She… resemble-s?', {"size": 18, "bold": True}),
    ('  resembles her mother."', {"size": 18, "bold": True}),
    ('  (잠깐 점검)', {"size": 14, "color": C_ASSOC}),
]
add_multiline(s, panel_x + Inches(0.3), panel_y + Inches(0.7),
              panel_w - Inches(0.6), Inches(3.6), ex_lines, line_space=1.2)

add_textbox(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.4),
            "Anderson (2000), Ch.9, p.319–322 — proceduralization, knowledge compilation",
            size=11, color=C_MUTED)

add_footer(s); add_page_number(s, 9, TOTAL)


# ============= SLIDE 10 — Stage 3 Autonomous =============
s = new_slide()
add_box(s, Inches(0.7), Inches(0.4), Inches(1.4), Inches(0.55), fill=C_AUTO, rounded=True)
add_textbox(s, Inches(0.7), Inches(0.45), Inches(1.4), Inches(0.5),
            "STAGE 3", size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER)
add_textbox(s, Inches(2.3), Inches(0.4), Inches(11.0), Inches(0.7),
            "Autonomous Stage", size=32, bold=True, color=C_TITLE)
add_textbox(s, Inches(2.3), Inches(1.05), Inches(11.0), Inches(0.4),
            "자율 단계  ·  Procedural memory 중심", size=18, color=C_MUTED)
add_box(s, Inches(0.7), Inches(1.55), Inches(1.2), Inches(0.04), fill=C_AUTO)

key_lines = [
    ("• 의식적으로 규칙을 떠올리지 않아도 정확", {}),
    ("• 자동화된 처리 절차·패턴으로 저장됨", {}),
    ("• 주의 자원이 의미·화용 처리로 이동", {}),
    ("• 빠르고 효율적인 산출이 가능", {"color": C_MUTED}),
]
add_multiline(s, Inches(0.9), Inches(1.95), Inches(7.0), Inches(3.5),
              key_lines, size=20, line_space=1.5)

# Right panel — example
panel_x = Inches(8.2); panel_y = Inches(1.95); panel_w = Inches(4.2); panel_h = Inches(4.5)
add_box(s, panel_x, panel_y, panel_w, panel_h, fill=RGBColor(0xF2, 0xF2, 0xF2),
        line=C_AUTO, line_w=1.5, rounded=True)
add_textbox(s, panel_x, panel_y + Inches(0.2), panel_w, Inches(0.4),
            "자연스러운 발화", size=12, bold=True, color=C_AUTO, align=PP_ALIGN.CENTER)
ex_lines = [
    ('잡담 중 ―', {"size": 14, "color": C_MUTED}),
    ('"She doesn\'t seem to', {"size": 18, "bold": True}),
    ('care, does she?"', {"size": 18, "bold": True}),
    ("", {"size": 6}),
    ('부가의문문 자동 산출', {"size": 14, "color": C_AUTO}),
    ('의식 없음 ✓', {"size": 14, "color": C_AUTO}),
    ("", {"size": 8}),
    ('회고적 인지 →', {"size": 14, "color": C_MUTED}),
    ('"내가 방금 부가의문문', {"size": 14}),
    ('을 썼지?"', {"size": 14}),
]
add_multiline(s, panel_x + Inches(0.3), panel_y + Inches(0.7),
              panel_w - Inches(0.6), Inches(3.7), ex_lines, line_space=1.2)

add_textbox(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.4),
            "Anderson (2000), Ch.9, p.325–326 — “requires less attention but is harder to interrupt”",
            size=11, color=C_MUTED)

add_footer(s); add_page_number(s, 10, TOTAL)


# ============= SLIDE 11 — Three Stages Integrated =============
s = new_slide()
add_title_bar(s, "Three Stages — Integrated View",
              "단계 간 변화를 한눈에")

# horizontal flow — 3 stage pills with arrows between
positions = [Inches(2.6), Inches(6.65), Inches(10.7)]
labels = [("Cognitive", C_COG), ("Associative", C_ASSOC), ("Autonomous", C_AUTO)]
pill_w = Inches(2.2); pill_h = Inches(1.0); cy = Inches(2.6)
for i, (cx, (lab, col)) in enumerate(zip(positions, labels)):
    px = cx - pill_w / 2
    add_box(s, px, cy, pill_w, pill_h, fill=col, line=col, rounded=True)
    add_textbox(s, px, cy + Inches(0.1), pill_w, Inches(0.4),
                f"Stage {i + 1}", size=12, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    add_textbox(s, px, cy + Inches(0.4), pill_w, Inches(0.55),
                lab, size=20, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

# arrows between pills
ay = cy + pill_h / 2
add_arrow(s, Inches(3.8), ay, Inches(5.45), ay, color=C_ACCENT, weight=4)
add_arrow(s, Inches(7.85), ay, Inches(9.5), ay, color=C_ACCENT, weight=4)

# Comparison table below
table_y = Inches(4.4)
headers = ["Dimension", "Cognitive", "Associative", "Autonomous"]
rows = [
    ["Memory type", "Declarative", "Compiling", "Procedural"],
    ["Speed", "Slow", "Faster", "Fast"],
    ["Errors", "Many", "Fewer", "Few"],
    ["Attention", "High", "Moderate", "Low"],
    ["Verbalizable?", "Yes", "Partly", "No"],
]
col_widths = [Inches(2.6), Inches(3.0), Inches(3.0), Inches(3.0)]
col_x = [Inches(0.85)]
for w in col_widths[:-1]:
    col_x.append(col_x[-1] + w)

# header row
hh = Inches(0.5)
for cx, w, txt in zip(col_x, col_widths, headers):
    add_box(s, cx, table_y, w, hh, fill=C_TITLE, line=C_TITLE, line_w=0.5)
    add_textbox(s, cx, table_y + Inches(0.05), w, Inches(0.4),
                txt, size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                align=PP_ALIGN.CENTER)

row_h = Inches(0.36)
stage_colors = [C_TEXT, C_COG, C_ASSOC, C_AUTO]
for r, row in enumerate(rows):
    ry = table_y + hh + r * row_h
    bg = RGBColor(0xF7, 0xF7, 0xF7) if r % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
    for j, (cx, w, txt) in enumerate(zip(col_x, col_widths, row)):
        add_box(s, cx, ry, w, row_h, fill=bg, line=C_LINE, line_w=0.5)
        col = stage_colors[j] if j > 0 else C_TEXT
        bold = (j > 0)
        add_textbox(s, cx, ry + Inches(0.03), w, Inches(0.3),
                    txt, size=13, bold=bold, color=col, align=PP_ALIGN.CENTER)

add_footer(s); add_page_number(s, 11, TOTAL)


# ============= SLIDE 12 — Production Rules =============
s = new_slide()
add_title_bar(s, "Production Rules", "If-Then  ·  Procedural knowledge의 기본 단위")

add_textbox(s, Inches(0.9), Inches(1.95), Inches(11.5), Inches(0.6),
            "“Production rules are condition-action pairs.”  — Anderson (2000), Ch.9",
            size=18, color=C_MUTED)

# 3 example boxes (escalating difficulty)
examples = [
    ("Simple", C_AUTO,
     "주어 = 3rd-person singular  AND  Tense = present",
     "Add  -s  to verb stem"),
    ("Intermediate", C_ASSOC,
     "주절 시제 = past  AND  종속절이 reported speech",
     "Backshift tense  (am → was, have → had)"),
    ("Complex", C_COG,
     "조건절 = 현재 사실 반대  AND  화자 = 가정·후회",
     "If + S + were  /  S + would-could-might + 동사원형"),
]

for i, (lev, col, cond, act) in enumerate(examples):
    y = Inches(2.8 + i * 1.45)
    # level badge
    add_box(s, Inches(0.9), y, Inches(1.4), Inches(1.2), fill=col, rounded=True)
    add_textbox(s, Inches(0.9), y + Inches(0.4), Inches(1.4), Inches(0.5),
                lev, size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                align=PP_ALIGN.CENTER)
    # IF box
    add_box(s, Inches(2.5), y, Inches(5.6), Inches(1.2),
            fill=C_LIGHT, line=C_ACCENT, line_w=1.0, rounded=True)
    add_textbox(s, Inches(2.65), y + Inches(0.1), Inches(0.5), Inches(0.4),
                "IF", size=14, bold=True, color=C_ACCENT)
    add_textbox(s, Inches(2.65), y + Inches(0.45), Inches(5.3), Inches(0.7),
                cond, size=15, color=C_TEXT)
    # arrow
    add_arrow(s, Inches(8.2), y + Inches(0.6), Inches(8.55), y + Inches(0.6),
              color=C_ACCENT, weight=3)
    # THEN box
    add_box(s, Inches(8.65), y, Inches(3.7), Inches(1.2),
            fill=RGBColor(0xF2, 0xF2, 0xF2), line=C_PROC, line_w=1.0, rounded=True)
    add_textbox(s, Inches(8.8), y + Inches(0.1), Inches(0.8), Inches(0.4),
                "THEN", size=14, bold=True, color=C_PROC)
    add_textbox(s, Inches(8.8), y + Inches(0.45), Inches(3.4), Inches(0.7),
                act, size=15, color=C_TEXT)

add_footer(s); add_page_number(s, 12, TOTAL)


# ============= SLIDE 13 — Knowledge Compilation =============
s = new_slide()
add_title_bar(s, "Knowledge Compilation",
              "Declarative → Procedural을 잇는 메커니즘")

# Two sub-mechanisms
sub_y = Inches(2.0)
# Composition
add_box(s, Inches(0.9), sub_y, Inches(5.7), Inches(2.2),
        fill=C_LIGHT, line=C_ACCENT, line_w=1.5, rounded=True)
add_textbox(s, Inches(0.9), sub_y + Inches(0.15), Inches(5.7), Inches(0.4),
            "1.  COMPOSITION", size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.9), sub_y + Inches(0.55), Inches(5.7), Inches(0.5),
            "여러 규칙을 하나로 합성", size=20, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(1.2), sub_y + Inches(1.15), Inches(5.1), Inches(1.0),
            "Rule A  +  Rule B  +  Rule C\n→  단일 production 으로 압축\n→  실행 속도 ↑",
            size=16, color=C_TEXT, align=PP_ALIGN.CENTER)

# Proceduralization
add_box(s, Inches(6.8), sub_y, Inches(5.7), Inches(2.2),
        fill=RGBColor(0xF2, 0xF2, 0xF2), line=C_PROC, line_w=1.5, rounded=True)
add_textbox(s, Inches(6.8), sub_y + Inches(0.15), Inches(5.7), Inches(0.4),
            "2.  PROCEDURALIZATION", size=14, bold=True, color=C_PROC,
            align=PP_ALIGN.CENTER)
add_textbox(s, Inches(6.8), sub_y + Inches(0.55), Inches(5.7), Inches(0.5),
            "선언적 사실을 행동에 박아넣기",
            size=20, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(7.1), sub_y + Inches(1.15), Inches(5.1), Inches(1.0),
            'Declarative info → IF-THEN 안에 내장\n작업기억 부담 ↓\n"verbal mediation drops out"',
            size=16, color=C_TEXT, align=PP_ALIGN.CENTER)

# Outcome strip
out_y = Inches(4.6)
add_box(s, Inches(0.9), out_y, Inches(11.6), Inches(0.7),
        fill=C_TITLE, line=C_TITLE, line_w=1.0, rounded=True)
add_textbox(s, Inches(0.9), out_y + Inches(0.18), Inches(11.6), Inches(0.4),
            "RESULT  ·  말로 되뇌이지 않아도 자동으로 실행되는 절차로 변환",
            size=18, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER)

# Caveat
cav_y = Inches(5.7)
add_box(s, Inches(0.9), cav_y, Inches(11.6), Inches(1.1),
        fill=RGBColor(0xF7, 0xF7, 0xF7), line=C_LINE, line_w=1.0, rounded=True)
add_textbox(s, Inches(1.1), cav_y + Inches(0.1), Inches(3.0), Inches(0.4),
            "⚠  IMPORTANT CAVEAT", size=12, bold=True, color=C_COG)
add_textbox(s, Inches(1.1), cav_y + Inches(0.45), Inches(11.0), Inches(0.6),
            "Skill-specific  —  연습한 방향(생성 vs 이해)으로만 자동화됨  (DeKeyser, 1997)",
            size=18, color=C_TEXT)

add_footer(s); add_page_number(s, 13, TOTAL)


# ============= SLIDE 14 — Power Law of Practice =============
s = new_slide()
add_title_bar(s, "Power Law of Practice",
              "연습이 자동화로 이어지는 수학적 패턴")

# generate curve image
curve_path = ASSET_DIR / "power_law.png"
make_power_curve(curve_path)
s.shapes.add_picture(str(curve_path), Inches(0.7), Inches(1.9), height=Inches(4.4))

# right side notes
nx = Inches(8.5); ny = Inches(2.0)
add_textbox(s, nx, ny, Inches(4.3), Inches(0.5),
            "RT  =  a · N⁻ᵇ", size=24, bold=True, color=C_TITLE)
add_textbox(s, nx, ny + Inches(0.55), Inches(4.3), Inches(0.45),
            "Newell & Rosenbloom (1981)", size=14, color=C_MUTED)

notes = [
    ("• 연습 횟수가 늘수록", {"size": 18}),
    ("  반응 시간이 멱함수로 감소", {"size": 18}),
    ("", {"size": 6}),
    ("• 초반 큰 향상 → 후반 점근선", {"size": 18}),
    ("  (negative acceleration)", {"size": 14, "color": C_MUTED}),
    ("", {"size": 6}),
    ("• 자동화는 ‘얻어지는 것’", {"size": 18}),
    ("  Automaticity is earned.", {"size": 14, "color": C_MUTED}),
]
add_multiline(s, nx, ny + Inches(1.2), Inches(4.3), Inches(4.0),
              notes, line_space=1.25)

add_textbox(s, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.4),
            "Logan (1988): 자동화 = 단일단계 인출(direct-access retrieval)로의 전환  —  경쟁가설",
            size=11, color=C_MUTED)

add_footer(s); add_page_number(s, 14, TOTAL)


# ============= SLIDE 15 — DeKeyser SLA =============
s = new_slide()
add_title_bar(s, "DeKeyser — Skill Acquisition in SLA",
              "Anderson의 인지이론을 외국어 학습에 적용")

add_textbox(s, Inches(0.9), Inches(1.95), Inches(11.5), Inches(0.6),
            "L2 문법 학습은 다른 모든 기능 학습과 동일한 경로를 따른다.",
            size=20, color=C_MUTED)

# 3-step path with arrows (vertical or horizontal)
y = Inches(2.9)
steps = [
    ("Declarative", "규칙·예문을\n명시적으로 학습", C_DECL),
    ("Proceduralization", "통제된 연습으로\n행동에 적용", C_ASSOC),
    ("Automatization", "Power-law 곡선\n반응시간↓ 오류↓", C_AUTO),
]
bw = Inches(3.5); bh = Inches(2.2); gap = Inches(0.4)
total = bw * 3 + gap * 2
sx = (SLIDE_W - total) / 2
for i, (head, body, col) in enumerate(steps):
    x = sx + (bw + gap) * i
    add_box(s, x, y, bw, bh, fill=C_BG, line=col, line_w=2.0, rounded=True)
    add_box(s, x, y, bw, Inches(0.5), fill=col, line=col, line_w=2.0, rounded=True)
    add_textbox(s, x, y + Inches(0.05), bw, Inches(0.4),
                head, size=16, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                align=PP_ALIGN.CENTER)
    add_textbox(s, x + Inches(0.2), y + Inches(0.75), bw - Inches(0.4),
                Inches(1.3), body, size=18, color=C_TEXT, align=PP_ALIGN.CENTER)
    if i < 2:
        add_arrow(s, x + bw + Inches(0.05), y + bh / 2,
                  x + bw + gap - Inches(0.05), y + bh / 2,
                  color=C_ACCENT, weight=3)

# Bottom finding
fy = Inches(5.4)
add_box(s, Inches(0.9), fy, Inches(11.6), Inches(1.4),
        fill=C_LIGHT, line=C_ACCENT, line_w=1.5, rounded=True)
add_textbox(s, Inches(1.1), fy + Inches(0.15), Inches(3.5), Inches(0.4),
            "KEY FINDING  ·  DeKeyser (1997)", size=12, bold=True, color=C_ACCENT)
add_textbox(s, Inches(1.1), fy + Inches(0.55), Inches(11.2), Inches(0.5),
            "한 학기 문법 연습  →  반응시간·오류율이 power-law 곡선",
            size=18, bold=True, color=C_TEXT)
add_textbox(s, Inches(1.1), fy + Inches(0.95), Inches(11.2), Inches(0.4),
            "“Practice the way you'll be tested.”  —  생성을 연습하면 생성만, 이해를 연습하면 이해만 자동화",
            size=14, color=C_MUTED)

add_footer(s); add_page_number(s, 15, TOTAL)


# ============= SLIDE 16 — L2 Grammar Examples =============
s = new_slide()
add_title_bar(s, "L2 Grammar Examples by Stage",
              "한국인 학습자의 영어 문법 발달")

# table
table_y = Inches(1.95)
headers = ["Stage", "Inner process", "Example utterance", "Speed"]
col_widths = [Inches(2.0), Inches(4.0), Inches(5.0), Inches(1.6)]
col_x = [Inches(0.85)]
for w in col_widths[:-1]:
    col_x.append(col_x[-1] + w)

hh = Inches(0.55)
for cx, w, txt in zip(col_x, col_widths, headers):
    add_box(s, cx, table_y, w, hh, fill=C_TITLE, line=C_TITLE, line_w=0.5)
    add_textbox(s, cx, table_y + Inches(0.1), w, Inches(0.4),
                txt, size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                align=PP_ALIGN.CENTER)

rows = [
    ("Cognitive", C_COG,
     "규칙을 의식적으로\n떠올리며 적용",
     '"3인칭 단수…-s 붙여야지"\n→  "She  goes  to  school."',
     "~5–7 s"),
    ("Associative", C_ASSOC,
     "익숙한 표현은 청크화\n새 어휘에선 멈칫",
     '"He goes ✓"\n"She… resemble-s? resembles."',
     "~1–2 s"),
    ("Autonomous", C_AUTO,
     "의식 없이 자동 산출\n주의는 의미·화용에",
     '"She doesn\'t seem to\ncare,  does  she?"',
     "<0.5 s"),
]

ry = table_y + hh
row_h = Inches(1.3)
for i, (stage, col, inner, utt, sp) in enumerate(rows):
    bg = RGBColor(0xF7, 0xF7, 0xF7) if i % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
    # row bg
    for cx, w in zip(col_x, col_widths):
        add_box(s, cx, ry, w, row_h, fill=bg, line=C_LINE, line_w=0.5)
    # stage badge
    add_box(s, col_x[0] + Inches(0.2), ry + Inches(0.45), col_widths[0] - Inches(0.4),
            Inches(0.65), fill=col, rounded=True)
    add_textbox(s, col_x[0], ry + Inches(0.55), col_widths[0], Inches(0.5),
                stage, size=15, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                align=PP_ALIGN.CENTER)
    # inner
    add_textbox(s, col_x[1] + Inches(0.15), ry + Inches(0.25),
                col_widths[1] - Inches(0.3), row_h - Inches(0.3),
                inner, size=15, color=C_TEXT, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    # utterance
    add_textbox(s, col_x[2] + Inches(0.15), ry + Inches(0.25),
                col_widths[2] - Inches(0.3), row_h - Inches(0.3),
                utt, size=15, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    # speed
    add_textbox(s, col_x[3], ry + Inches(0.25), col_widths[3], row_h - Inches(0.3),
                sp, size=16, bold=True, color=col, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    ry += row_h

add_footer(s); add_page_number(s, 16, TOTAL)


# ============= SLIDE 17 — Ullman DP Model =============
s = new_slide()
add_title_bar(s, "Ullman's Declarative / Procedural Model",
              "신경언어학적 기반  ·  뇌의 두 시스템이 언어를 처리한다")

# Two-side layout
ly = Inches(2.0); lh = Inches(3.4)
# Left — Declarative (lexicon)
add_box(s, Inches(0.9), ly, Inches(5.7), lh, fill=C_LIGHT, line=C_DECL,
        line_w=2.0, rounded=True)
add_textbox(s, Inches(0.9), ly + Inches(0.15), Inches(5.7), Inches(0.4),
            "DECLARATIVE", size=13, bold=True, color=C_DECL, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.9), ly + Inches(0.55), Inches(5.7), Inches(0.5),
            "Mental Lexicon", size=24, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)
left_lines = [
    ("· 측두엽 / 해마 (Temporal · Hippocampus)", {"size": 16}),
    ("· 단어, 불규칙형, 어휘적 사실", {"size": 16}),
    ("· “book, went, children”", {"size": 16, "bold": True, "color": C_DECL}),
    ("· 명시적 회상 가능", {"size": 16}),
]
add_multiline(s, Inches(1.3), ly + Inches(1.4), Inches(5.0), Inches(2.0),
              left_lines, line_space=1.5)

# Right — Procedural (grammar)
add_box(s, Inches(7.0), ly, Inches(5.5), lh, fill=RGBColor(0xF2, 0xF2, 0xF2),
        line=C_PROC, line_w=2.0, rounded=True)
add_textbox(s, Inches(7.0), ly + Inches(0.15), Inches(5.5), Inches(0.4),
            "PROCEDURAL", size=13, bold=True, color=C_PROC, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(7.0), ly + Inches(0.55), Inches(5.5), Inches(0.5),
            "Mental Grammar", size=24, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)
right_lines = [
    ("· 전두엽 + 기저핵 / 소뇌", {"size": 16}),
    ("· 규칙 기반 결합, 통사 처리", {"size": 16}),
    ("· walk → walked  (rule-based)", {"size": 16, "bold": True, "color": C_PROC}),
    ("· 무의식·자동", {"size": 16}),
]
add_multiline(s, Inches(7.4), ly + Inches(1.4), Inches(4.8), Inches(2.0),
              right_lines, line_space=1.5)

# Bottom — L1 vs L2 contrast
cy_box = Inches(5.55)
add_box(s, Inches(0.9), cy_box, Inches(11.6), Inches(1.45),
        fill=RGBColor(0xF7, 0xF7, 0xF7), line=C_LINE, line_w=1.0, rounded=True)
add_textbox(s, Inches(1.1), cy_box + Inches(0.12), Inches(3.0), Inches(0.4),
            "L1 vs L2  CONTRAST", size=12, bold=True, color=C_ACCENT)
contrast = [
    ("L1:  문법도 procedural로 처리  (자동, 무의식)", {"size": 15}),
    ("L2 초기:  문법조차 declarative에 의존  →  숙달과 함께 procedural로 점진적 이동",
     {"size": 15, "bold": True, "color": C_TEXT}),
]
add_multiline(s, Inches(1.1), cy_box + Inches(0.5), Inches(11.2), Inches(0.9),
              contrast, line_space=1.3)

add_footer(s); add_page_number(s, 17, TOTAL)


# ============= SLIDE 18 — Summary Diagram =============
s = new_slide()
add_title_bar(s, "Summary  ·  The Whole Picture",
              "외국어 문법 학습의 전체 흐름")

# vertical flow
y = Inches(1.9)
center_x = SLIDE_W / 2
# Top: input
add_box(s, center_x - Inches(3.0), y, Inches(6.0), Inches(0.7),
        fill=C_LIGHT, line=C_ACCENT, line_w=1.5, rounded=True)
add_textbox(s, center_x - Inches(3.0), y + Inches(0.15), Inches(6.0), Inches(0.5),
            "교사 설명 · 예문 · 명시적 규칙",
            size=18, bold=True, color=C_TITLE, align=PP_ALIGN.CENTER)

# stages diagram
sy = Inches(2.9)
stage_h = Inches(0.85)
gap_v = Inches(0.15)
stages_v = [
    ("STAGE 1  ·  Cognitive", "Declarative memory에 규칙·예문 저장", C_COG),
    ("STAGE 2  ·  Associative", "Knowledge compilation으로 IF-THEN화", C_ASSOC),
    ("STAGE 3  ·  Autonomous", "Procedural memory에서 자동 산출", C_AUTO),
]
for i, (head, body, col) in enumerate(stages_v):
    yy = sy + (stage_h + gap_v) * i
    add_box(s, center_x - Inches(4.0), yy, Inches(8.0), stage_h,
            fill=C_BG, line=col, line_w=2.0, rounded=True)
    # left badge
    add_box(s, center_x - Inches(4.0), yy, Inches(2.5), stage_h, fill=col, rounded=True)
    add_textbox(s, center_x - Inches(4.0), yy + Inches(0.22), Inches(2.5), Inches(0.5),
                head.split("·")[0].strip(), size=14, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    add_textbox(s, center_x - Inches(1.4), yy + Inches(0.22), Inches(5.4), Inches(0.5),
                body, size=16, color=C_TEXT)
    if i < 2:
        ax = center_x
        ay1 = yy + stage_h
        ay2 = yy + stage_h + gap_v
        add_arrow(s, ax, ay1, ax, ay2, color=C_ACCENT, weight=3)

# bottom: output
oy = sy + (stage_h + gap_v) * 3 + Inches(0.05)
add_arrow(s, center_x, sy + (stage_h + gap_v) * 3 - Inches(0.1),
          center_x, oy, color=C_ACCENT, weight=3)
add_box(s, center_x - Inches(3.5), oy + Inches(0.15), Inches(7.0), Inches(0.7),
        fill=C_TITLE, line=C_TITLE, line_w=1.5, rounded=True)
add_textbox(s, center_x - Inches(3.5), oy + Inches(0.3), Inches(7.0), Inches(0.5),
            "유창하고 자동화된 L2 문법 사용",
            size=18, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

add_footer(s); add_page_number(s, 18, TOTAL)


# ============= SLIDE 19 — Key Takeaways =============
s = new_slide()
add_title_bar(s, "Key Takeaways", "꼭 기억해야 할 다섯 가지")

takeaways = [
    ("Knowing  ≠  Using",
     "선언적으로 ‘아는 것’과 절차적으로 ‘쓰는 것’은 다르다"),
    ("From Rules to Reflexes",
     "Declarative → Procedural로 기억 시스템이 점진적으로 전이된다"),
    ("Practice Rewrites the Brain",
     "Knowledge compilation: IF-THEN production으로 재구조화"),
    ("Automaticity Is Earned",
     "Power-law 곡선 — 충분한 연습만이 자동화를 만든다"),
    ("Practice the Way You'll Be Tested",
     "Skill-specific — 연습한 방향으로만 자동화된다 (DeKeyser)"),
]
for i, (head, body) in enumerate(takeaways):
    y = Inches(1.95 + i * 1.05)
    # number circle
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), y, Inches(0.7), Inches(0.7))
    circ.fill.solid(); circ.fill.fore_color.rgb = C_ACCENT
    circ.line.color.rgb = C_ACCENT
    add_textbox(s, Inches(0.9), y + Inches(0.13), Inches(0.7), Inches(0.5),
                str(i + 1), size=20, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1.85), y - Inches(0.05), Inches(11.0), Inches(0.5),
                head, size=22, bold=True, color=C_TITLE)
    add_textbox(s, Inches(1.85), y + Inches(0.45), Inches(11.0), Inches(0.5),
                body, size=15, color=C_MUTED)

add_footer(s); add_page_number(s, 19, TOTAL)


# ============= SLIDE 20 — References =============
s = new_slide()
add_title_bar(s, "References", "주요 참고 문헌")

refs = [
    "Anderson, J. R. (1982). Acquisition of cognitive skill. Psychological Review, 89(4), 369–406.",
    "Anderson, J. R. (2000). Learning and Memory: An Integrated Approach (2nd ed.). Wiley.",
    "Fitts, P. M., & Posner, M. I. (1967). Human Performance. Brooks/Cole.",
    "Newell, A., & Rosenbloom, P. S. (1981). Mechanisms of skill acquisition and the law of practice.",
    "Logan, G. D. (1988). Toward an instance theory of automatization. Psychological Review, 95(4), 492–527.",
    "DeKeyser, R. M. (1997). Beyond explicit rule learning: Automatizing second language morphosyntax. SSLA, 19(2), 195–221.",
    "DeKeyser, R. M. (2007). Practice in a Second Language. Cambridge University Press.",
    "Ullman, M. T. (2001). The neural basis of lexicon and grammar in first and second language. Bilingualism, 4(2), 105–122.",
    "Ullman, M. T. (2004). Contributions of memory circuits to language: The declarative/procedural model. Cognition, 92(1–2), 231–270.",
    "Squire, L. R. (1992). Declarative and nondeclarative memory: Multiple brain systems supporting learning and memory. JoCN, 4(3), 232–243.",
]
y = Inches(1.95)
for r in refs:
    add_textbox(s, Inches(0.9), y, Inches(11.6), Inches(0.45),
                "·  " + r, size=13, color=C_TEXT)
    y += Inches(0.42)

# closing line
add_textbox(s, Inches(0.9), Inches(6.5), Inches(11.6), Inches(0.5),
            "Thank you.", size=22, bold=True, color=C_TITLE, align=PP_ALIGN.CENTER)

add_footer(s); add_page_number(s, 20, TOTAL)


# ============= Save =============
prs.save(PPT_PATH)
print(f"✓ Saved: {PPT_PATH}")
print(f"  slides: {len(prs.slides)}")
